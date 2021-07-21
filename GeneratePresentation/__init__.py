import logging

import azure.functions as func
import os

from azure.storage.blob import BlobClient, BlobSasPermissions, generate_blob_sas
from datetime import datetime, timedelta, timezone

from pptx import Presentation
from pptx.util import Pt
import requests # pulling data
from bs4 import BeautifulSoup # xml parsing

# RSS scraping function
# Based mostly on: https://github.com/mattdood/web_scraping_example/blob/master/scraping.py
def get_updates_rss(startDate, endDate):
    article_list = []

    try:
        # execute my request, parse the data using XML
        # parse using BS4
        r = requests.get(os.environ["UpdatesURL"])
        soup = BeautifulSoup(r.content, features='xml')

        # select only the "items" I want from the data
        updates = soup.findAll('item')

        # for each "item" I want, parse it into a list
        for a in updates:

            # Get publication date
            published = a.find('pubDate').text
            pubDate = datetime.strptime(a.find('pubDate').text, "%a, %d %b %Y %H:%M:%S Z")

            # only include items falling within our requested date range
            if (pubDate >= startDate and pubDate <= endDate):

                title = a.find('title').text
                link = a.find('link').text
            
                # basic parse to flag announcement types
                if "preview" in title.lower():
                     announcement_type = "preview"
                else:
                    announcement_type = "GA"

                # create an "article" object with the data
                # from each "item"
                article = {
                    'title': title,
                    'link': link,
                    'published': published,
                    'antype': announcement_type
                    }

                # append my "article_list" with each "article" object
                article_list.append(article)
        
        # after the loop, dump my saved objects into a .txt file
        return article_list
    except Exception as e:
        logging.exception("Couldn't scrape the Azure Updates RSS feed")

###
# Generate a section of the final PowerPoint
###
def generate_presentation_section(presentation, layout, articles, item_type):

    # Add first slide and slide notes
    slide = presentation.slides.add_slide(layout)
    slide_notes = slide.notes_slide
    shapes = slide.shapes

    slide_item_count = 0
    slide_count = 1

    for article in articles:

        # Each new slide requires first elements be added differently to the rest.
        if slide_item_count == 0:

            # Insert title for slide
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = item_type + " (" + str(slide_count) + ")"
            # Insert first bullet item
            tf = body_shape.text_frame
            tf.text = article["title"]
            tf.paragraphs[0].font.size = Pt(24)
            # Insert first slide note
            sltf = slide_notes.notes_text_frame
            sltf.text = "- " + article["link"] + " (" + article["published"] + ")"

        else:

            # Insert bullet point
            p = tf.add_paragraph()
            p.font.size = Pt(24)
            p.text = article["title"]
            # Insert slide note
            dotpoint = sltf.add_paragraph()
            dotpoint.text = "- " + article["link"] + " (" + article["published"] + ")"
        
        slide_item_count += 1

        # If we hit 5 items on a slide, create a new slide and reset item count
        if slide_item_count == 5:
            slide = presentation.slides.add_slide(layout)
            slide_notes = slide.notes_slide
            shapes = slide.shapes
            slide_item_count = 0
            slide_count += 1

###
# Upload generated file to Azure Storage and generate a SAS URL for it
###
def upload_file_to_storage(presenation_file):

    blob = BlobClient.from_connection_string(conn_str=os.environ["PowerPointAccountConnection"], container_name=os.environ["PowerPointContainer"], blob_name=presenation_file)

    with open(presenation_file, "rb") as data:
        blob.upload_blob(data)

    # Generate a SAS-protected URL for the item which will allow the caller to download the file for 1 hour.
    startTime = datetime.now(tz=timezone.utc)
    endTime = startTime + timedelta(hours=1)
    return "https://" + os.environ["PowerPointStorageAccount"] + ".blob.core.windows.net/" + os.environ["PowerPointContainer"] + "/" + presenation_file + "?" + generate_blob_sas(os.environ["PowerPointStorageAccount"],os.environ["PowerPointContainer"],blob_name=presenation_file,account_key=os.environ["PowerPointStorageKey"],permission=BlobSasPermissions(read=True),start=startTime,expiry=endTime)

#####
# Azure Function main entry point
#####
def main(req: func.HttpRequest) -> func.HttpResponse:

    blob_sas_url = ""
    message = ""
    http_status = 200

    try:

        # start date is required
        startParam = req.params.get('start')
        if not startParam:
            
            message = "Bad request: 'start' query parameter is required in format YYYY-MM-DD."
            http_status=400

        else:

            # end date is optional, so if not provided use today
            endParam = req.params.get('end')
            if not endParam:
                endParam = datetime.now("%Y-%m-%d")

            # add 1 day to end date so we include all of the day
            ending = datetime.strptime(endParam, "%Y-%m-%d")
            ending = ending + timedelta(days=1)
            starting = datetime.strptime(startParam, "%Y-%m-%d")

            updatelist = get_updates_rss(startDate=starting,endDate=ending)

            if len(updatelist) > 0:

                prs = Presentation()
                # Initialise default slide layout (bullets)
                bullet_slide_layout = prs.slide_layouts[1]

                preview_items = [item for item in updatelist if item["antype"] == "preview"]
                ga_items = [item for item in updatelist if item["antype"] == "GA"]

                generate_presentation_section(prs, bullet_slide_layout, preview_items, "Preview")
                generate_presentation_section(prs, bullet_slide_layout, ga_items, "GA")
            
                filename = datetime.strftime(datetime.now(),"%Y-%m-%d-%H-%M-%S") + "-AzureUpdates.pptx"

                prs.save(filename)

                blob_sas_url = upload_file_to_storage(filename)

                message = "File created and uploaded to storage. You can <a href='" + blob_sas_url + "'>download it</a> for the next 1 hour."
            else:
                message = "There are no updates for the specified period, so no PowerPoint has been generated.",
        
    except TypeError as te:
        logging.exception("Type error")
        message = "Check the format of your request and ensure you provide the 'start' query parameter in the format YYYY-MM-DD",
        http_status=400

    except ValueError:
        pass

    return func.HttpResponse(
            mimetype="text/html",
            body=message,
            status_code=http_status
    )